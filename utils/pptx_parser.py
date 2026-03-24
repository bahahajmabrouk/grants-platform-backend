from pptx import Presentation
from pathlib import Path


def extract_text_from_pptx(file_path: str) -> str:
    """
    Extrait tout le texte d'un fichier PPTX slide par slide.
    Inclut titres, corps de texte et notes du présentateur.
    """
    prs = Presentation(file_path)
    text_content = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)

        # Notes du présentateur (souvent très riches en info)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_texts.append(f"[Notes]: {notes}")

        if slide_texts:
            text_content.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))

    return "\n\n".join(text_content)


def get_pptx_metadata(file_path: str) -> dict:
    """Retourne les métadonnées du PPTX"""
    prs = Presentation(file_path)
    return {
        "slide_count": len(prs.slides),
        "file_size_kb": round(Path(file_path).stat().st_size / 1024, 1),
    }
